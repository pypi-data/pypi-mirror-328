"""
This is the main module for the kvell_extraction package.
"""

import argparse
import warnings
from pathlib import Path
from typing import Dict, List, Tuple, Union

import cv2
import filetype
import fitz
import numpy as np
from spire.doc import Document
from spire.doc.common import *
import polars as pl
from pptx import Presentation

from .utils import import_package


class PDFExtracter:
    """
    Class for extracting text from PDF files using OCR.
    """

    def __init__(self, dpi=200, **ocr_kwargs):
        """
        Initialize the PDFExtracter.
        """
        self.dpi = dpi

        ocr_engine = import_package("rapidocr_onnxruntime")

        if ocr_engine is None:  # Fix: Remove extra indentation
            raise ModuleNotFoundError(
                "Can't find the rapidocr_onnxruntime.\n Please pip install rapidocr_onnxruntime to run the code."
            )

        self.text_sys = ocr_engine.RapidOCR(**ocr_kwargs)
        self.empyt_list = []

    def __call__(
        self,
        content: Union[str, Path, bytes],
        force_ocr: bool = False,
    ) -> List[List[Union[str, str, str]]]:
        """
        Extract text from the input content.
        """
        try:
            file_type = self.which_type(content)
        except (FileExistsError, TypeError) as e:
            raise PDFExtracterError("The input content is empty.") from e

        supported_images = ["jpg", "jpeg", "png", "bmp"]
        if file_type not in ["pdf"] + supported_images:
            raise PDFExtracterError(
                f"The file type must be PDF or one of {supported_images}"
            )

        try:
            if file_type == "pdf":
                return self._handle_pdf(content, force_ocr)
            else:
                return self._handle_image(content)
        except PDFExtracterError as e:
            warnings.warn(str(e))
            return self.empyt_list

    def _handle_pdf(
        self,
        content: Union[str, Path, bytes],
        force_ocr: bool = False,
    ) -> List[List[Union[str, str, str]]]:
        """
        Extract text from a PDF file.
        """
        pdf_data = self.load_pdf(content)
        txts_dict, need_ocr_idxs = self.extract_texts(pdf_data, force_ocr)
        page_img_dict = self.read_pdf_with_image(pdf_data, need_ocr_idxs)
        ocr_res_dict = self.get_ocr_res(page_img_dict)
        return self.merge_direct_ocr(txts_dict, ocr_res_dict)

    def _handle_image(
        self,
        content: Union[str, Path, bytes],
    ) -> List[List[Union[str, str, str]]]:
        """
        Extract text from an image file.
        """
        if isinstance(content, (str, Path)):
            img_path = str(content)
            preds, _ = self.text_sys(img_path)
        else:
            # For bytes input, convert to numpy array first
            nparr = np.frombuffer(content, np.uint8)
            img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)
            preds, _ = self.text_sys(img)

        if preds:
            _, rec_res, _ = list(zip(*preds))
            text = "\n".join(rec_res)
            return [["0", text, "1.0"]]
        return self.empyt_list

    @staticmethod
    def load_pdf(pdf_content: Union[str, Path, bytes]) -> bytes:
        """
        Load the PDF content.
        """
        if isinstance(pdf_content, (str, Path)):
            if not Path(pdf_content).exists():
                raise PDFExtracterError(f"{pdf_content} does not exist.")

            with open(pdf_content, "rb") as f:
                data = f.read()
            return data

        if isinstance(pdf_content, bytes):
            return pdf_content

        raise PDFExtracterError(f"{type(pdf_content)} is not in [str, Path, bytes].")

    def extract_texts(self, pdf_data: bytes, force_ocr: bool) -> Tuple[Dict, List]:
        """
        Extract texts from the PDF content.
        """
        texts, need_ocr_idxs = {}, []
        with fitz.open(stream=pdf_data) as doc:
            for i, page in enumerate(doc):
                if force_ocr:
                    need_ocr_idxs.append(i)
                    continue

                text = page.get_text("text", sort=True)
                if text:
                    texts[str(i)] = text
                else:
                    need_ocr_idxs.append(i)
        return texts, need_ocr_idxs

    def read_pdf_with_image(self, pdf_data: bytes, need_ocr_idxs: List) -> Dict:
        """
        Read the PDF content with images.
        """

        def convert_img(page):
            pix = page.get_pixmap(dpi=self.dpi)
            img = np.frombuffer(pix.samples, dtype=np.uint8)
            img = img.reshape([pix.h, pix.w, pix.n])
            img = cv2.cvtColor(img, cv2.COLOR_BGR2RGB)
            return img

        with fitz.open(stream=pdf_data) as doc:
            page_img_dict = {k: convert_img(doc[k]) for k in need_ocr_idxs}
        return page_img_dict

    def get_ocr_res(self, page_img_dict: Dict) -> Dict:
        """
        Get the OCR results.
        """
        ocr_res = {}
        for k, v in page_img_dict.items():
            preds, _ = self.text_sys(v)
            if preds:
                _, rec_res, _ = list(zip(*preds))
                ocr_res[str(k)] = "\n".join(rec_res)
        return ocr_res

    def merge_direct_ocr(self, txts_dict, ocr_res_dict):
        """
        Merge the direct OCR results.
        """
        final_result = {**txts_dict, **ocr_res_dict}
        final_result = dict(sorted(final_result.items(), key=lambda x: int(x[0])))
        final_result = [[k, v, "1.0"] for k, v in final_result.items()]
        return final_result

    @staticmethod
    def which_type(content: Union[bytes, str, Path]) -> str:
        """
        Determine the type of the input content.
        First tries using filetype library, then falls back to file extension.
        """
        if isinstance(content, (str, Path)) and not Path(content).exists():
            raise FileExistsError(f"{content} does not exist.")

        # First try using filetype
        kind = filetype.guess(content)
        if kind is not None:
            return kind.extension

        # If filetype fails, try using file extension
        if isinstance(content, (str, Path)):
            ext = Path(content).suffix.lower()[1:]  # Remove the dot
            if ext:
                return ext

        raise TypeError(f"The type of {content} does not support.")


class PDFExtracterError(Exception):
    """Exception for PDFExtracter errors"""

    pass


class ImageExtracter(PDFExtracter):
    """Class for extracting text from images using OCR"""

    def __call__(
        self,
        content: Union[str, Path, bytes],
        force_ocr: bool = False,  # kept for compatibility
    ) -> List[List[Union[str, str, str]]]:
        try:
            file_type = self.which_type(content)
        except (FileExistsError, TypeError) as e:
            raise PDFExtracterError("The input content is empty.") from e

        supported_images = ["jpg", "jpeg", "png", "bmp"]
        if file_type not in supported_images:
            raise PDFExtracterError(f"The file type must be one of {supported_images}")

        try:
            return self._handle_image(content)
        except PDFExtracterError as e:
            warnings.warn(str(e))
            return self.empyt_list


class DocExtracter(PDFExtracter):
    """Class for extracting text from Word documents using Spire.Doc"""

    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        self.supported_formats = ["doc", "docx", "docm", "dot", "dotx", "dotm"]

    def __call__(
        self,
        content: Union[str, Path, bytes],
        force_ocr: bool = False,  # kept for compatibility
    ) -> List[List[Union[str, str, str]]]:
        try:
            file_type = self.which_type(content)
        except (FileExistsError, TypeError) as e:
            raise PDFExtracterError("The input content is empty.") from e

        if file_type not in self.supported_formats:
            raise PDFExtracterError(
                f"The file type must be one of {self.supported_formats}"
            )

        try:
            return self._handle_doc(content)
        except Exception as e:
            warnings.warn(str(e))
            return self.empyt_list

    def _handle_doc(
        self,
        content: Union[str, Path, bytes],
    ) -> List[List[Union[str, str, str]]]:
        document = Document()

        if isinstance(content, (str, Path)):
            document.LoadFromFile(str(content))
        elif isinstance(content, bytes):
            # Load from memory stream
            from io import BytesIO

            stream = BytesIO(content)
            document.LoadFromStream(stream)

        text = document.GetText()
        return [["0", text, "1.0"]]  # Return as single page with confidence 1.0


class ExcelExtracter(PDFExtracter):
    """Class for extracting text from Excel files using polars"""

    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        self.supported_formats = ['xlsx', 'xls']

    def __call__(
        self,
        content: Union[str, Path, bytes],
        force_ocr: bool = False,  # kept for compatibility
    ) -> List[List[Union[str, str, str]]]:
        try:
            # For Excel files, check extension first since filetype may not detect .xls properly
            if isinstance(content, (str, Path)):
                ext = Path(content).suffix.lower()[1:]  # Remove the dot
                if ext in self.supported_formats:
                    file_type = ext
                else:
                    # If extension not supported, try filetype detection as fallback
                    file_type = self.which_type(content)
            else:
                file_type = self.which_type(content)
                
        except (FileExistsError, TypeError) as e:
            raise PDFExtracterError("The input content is empty.") from e

        if file_type not in self.supported_formats:
            raise PDFExtracterError(
                f"The file type must be one of {self.supported_formats}"
            )

        try:
            return self._handle_excel(content)
        except Exception as e:
            warnings.warn(str(e))
            return self.empyt_list

    def _handle_excel(
        self,
        content: Union[str, Path, bytes],
    ) -> List[List[Union[str, str, str]]]:
        if isinstance(content, (str, Path)):
            df = pl.read_excel(str(content))
        elif isinstance(content, bytes):
            # For bytes input, save to temporary file first
            import tempfile

            with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            df = pl.read_excel(tmp_path)
            import os

            os.unlink(tmp_path)  # Clean up temp file

        # Convert all columns to strings
        df_text = df.select([pl.col("*").cast(pl.Utf8)])

        # Convert to text with tab separation
        rows_text = []
        for row in df_text.iter_rows():
            row_text = "\t".join(map(str, row))
            rows_text.append(row_text)

        # Join all rows with newlines
        full_text = "\n".join(rows_text)

        return [["0", full_text, "1.0"]]  # Return as single page with confidence 1.0


class PresentationExtracter(PDFExtracter):
    """Class for extracting text from PowerPoint presentations using python-pptx"""

    def __init__(self, **kwargs):
        super().__init__(**kwargs)
        self.supported_formats = ["pptx", "potx"]

    def __call__(
        self,
        content: Union[str, Path, bytes],
        force_ocr: bool = False,  # kept for compatibility
    ) -> List[List[Union[str, str, str]]]:
        try:
            file_type = self.which_type(content)
        except (FileExistsError, TypeError) as e:
            raise PDFExtracterError("The input content is empty.") from e

        if file_type not in self.supported_formats:
            raise PDFExtracterError(
                f"The file type must be one of {self.supported_formats}"
            )

        try:
            return self._handle_presentation(content)
        except Exception as e:
            warnings.warn(str(e))
            return self.empyt_list

    def _handle_presentation(
        self,
        content: Union[str, Path, bytes],
    ) -> List[List[Union[str, str, str]]]:
        if isinstance(content, (str, Path)):
            prs = Presentation(str(content))
        elif isinstance(content, bytes):
            # For bytes input, save to temporary file first
            import tempfile

            with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
                tmp.write(content)
                tmp_path = tmp.name
            prs = Presentation(tmp_path)
            import os

            os.unlink(tmp_path)  # Clean up temp file

        # Extract text from all slides
        slides_text = []
        for slide_number, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:  # Only include if there's actual text
                        slide_texts.append(text)

            if slide_texts:  # Only include slides with text
                slides_text.append(
                    [
                        str(slide_number),  # page number
                        "\n".join(slide_texts),  # text content
                        "1.0",  # confidence
                    ]
                )

        return slides_text if slides_text else [["0", "", "1.0"]]


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument(
        "-path",
        "--file_path",
        type=str,
        help="File path, supports PDF, Word documents, Excel files, and images",
    )
    parser.add_argument(
        "-f",
        "--force_ocr",
        action="store_true",
        default=False,
        help="Whether to use ocr for all pages (PDF only).",
    )
    args = parser.parse_args()

    # Determine file type and use appropriate extracter
    try:
        file_type = PDFExtracter.which_type(args.file_path)
        if file_type == "pdf":
            extracter = PDFExtracter()
        elif file_type in ["doc", "docx", "docm", "dot", "dotx", "dotm"]:
            extracter = DocExtracter()
        elif file_type in ["xlsx", "xls"]:
            extracter = ExcelExtracter()
        elif file_type in ["pptx", "potx"]:
            extracter = PresentationExtracter()
        elif file_type in ["jpg", "jpeg", "png", "bmp"]:
            extracter = ImageExtracter()
        else:
            raise PDFExtracterError(f"Unsupported file type: {file_type}")

        result = extracter(args.file_path)
        print(result)
    except Exception as e:
        print(f"Error: {str(e)}")


if __name__ == "__main__":
    main()
