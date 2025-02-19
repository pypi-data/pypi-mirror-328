from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

class PPTXEditor:
    def __init__(self, file_path):
        """Initialize with a PowerPoint file and set current slide index to 0."""
        self.file_path = file_path
        self.prs = Presentation(file_path)
        self.current_slide_index = 0  # Default starting slide

    def next_slide(self):
        """Move to the next slide if possible."""
        if self.current_slide_index < len(self.prs.slides) - 1:
            self.current_slide_index += 1
        else:
            print("Already on the last slide.")

    def previous_slide(self):
        """Move to the previous slide if possible."""
        if self.current_slide_index > 0:
            self.current_slide_index -= 1
        else:
            print("Already on the first slide.")

    def go_to_slide(self, slide_index):
        """Jump to a specific slide index if within bounds."""
        if 0 <= slide_index < len(self.prs.slides):
            self.current_slide_index = slide_index
        else:
            raise ValueError(f"Slide index must be between 0 and {len(self.prs.slides) - 1}")

    def add_red_dot(self, x=1, y=1, size=0.3):
        # """Add a red dot to the current slide."""
        # slide = self.prs.slides[self.current_slide_index]
        # left = Inches(x)
        # top = Inches(y)
        # width = Inches(size)
        # height = Inches(size)

        # shape = slide.shapes.add_shape(
        #     1, left, top, width, height  # 1 is for ellipse shape
        # )
        # shape.fill.solid()
        # shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color
        # shape.line.fill.background()
        """Add a red dot to the current slide at (x, y) inches."""
        slide = self.prs.slides[self.current_slide_index]
        left = Inches(x)
        top = Inches(y)
        diameter = Inches(size)  # Ensure width and height are equal for a circle

        shape = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, left, top, diameter, diameter  # OVAL makes a circle when width = height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color
        shape.line.fill.background()  # Remove border

    def save_presentation(self, output_path):
        """Save the modified PowerPoint file."""
        self.prs.save(output_path)

    def get_current_slide_index(self):
        """Return the current slide index."""
        return self.current_slide_index

    def get_total_slides(self):
        """Return the total number of slides in the presentation."""
        return len(self.prs.slides)
